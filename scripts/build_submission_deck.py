"""Build the 5-slide submission pitch deck.

Required by OpenClaw 2026 rules:
  - Max 5 slides
  - PDF format
  - Order: Problem · Solution · Architecture · Features+Stack · Future/Impact
  - Team and member names on the title/first slide (no separate intro slide)

Output:
  /home/ubuntu/ledgersoul/OpenClaw2026_TokenBurner_LedgerSoul.pptx
  /home/ubuntu/ledgersoul/OpenClaw2026_TokenBurner_LedgerSoul.pdf
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---- Theme ---------------------------------------------------------------
BG = RGBColor(0x07, 0x0B, 0x14)
CYAN = RGBColor(0x4D, 0xE6, 0xF0)
PURPLE = RGBColor(0xA9, 0x7C, 0xFF)
CARD = RGBColor(0x18, 0x24, 0x40)
CARD_EDGE = RGBColor(0x22, 0x31, 0x55)
ROW_BG = RGBColor(0x0D, 0x14, 0x28)
TEXT = RGBColor(0xE8, 0xEC, 0xF6)
DIM = RGBColor(0x9A, 0xA6, 0xC2)
GREEN = RGBColor(0x69, 0xE6, 0xA3)
RED = RGBColor(0xFF, 0x8A, 0x8A)
GOLD = RGBColor(0xFF, 0xC8, 0x6B)

INTER = "Inter"
MONO = "JetBrains Mono"

TEAM = "TokenBurner"
PROJECT = "LedgerSoul"


# ---- Helpers -------------------------------------------------------------

def add_rect(slide, x, y, w, h, fill, line=None, rounded=False, line_pt=0.75):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_pt)
    s.shadow.inherit = False
    if rounded:
        s.adjustments[0] = 0.10
    return s


def add_text(slide, x, y, w, h, text, *, font=INTER, size=14, bold=False,
             color=TEXT, align="left", anchor="top", italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def chrome(slide, eyebrow, title, subtitle, page_num, total=5):
    add_rect(slide, 0, 0, 13.33, 7.5, BG)
    add_rect(slide, 0.60, 0.55, 0.18, 0.18, CYAN)
    add_text(slide, 0.85, 0.50, 8.0, 0.30, eyebrow,
             font=MONO, size=10, bold=True, color=CYAN)
    add_text(slide, 0.60, 0.85, 12.10, 1.05, title,
             font=INTER, size=30, bold=True, color=TEXT)
    if subtitle:
        add_text(slide, 0.60, 1.95, 11.50, 0.40, subtitle,
                 font=INTER, size=14, color=DIM)
    add_text(slide, 0.60, 7.05, 9.5, 0.30,
             f"OPENCLAW 2026  ·  TEAM {TEAM.upper()}  ·  {PROJECT.upper()}",
             font=MONO, size=8, bold=True, color=DIM)
    add_text(slide, 11.60, 7.05, 1.20, 0.30,
             f"{page_num:02d} / {total:02d}",
             font=MONO, size=8, bold=True, color=DIM, align="right")


# ---- Slide 1: Problem ----------------------------------------------------

def slide_problem(slide):
    chrome(
        slide,
        eyebrow="01  ·  PROBLEM",
        title="Indonesian merchants on DOKU\nare drowning in payment ops.",
        subtitle="Every failed payment, refund, and suspicious charge becomes a manual firefight across dashboards, Slack, and inboxes.",
        page_num=1,
    )

    # Top-left: team strip
    add_rect(slide, 0.60, 2.50, 12.10, 0.45, CARD, line=CARD_EDGE, rounded=True)
    add_text(slide, 0.95, 2.50, 11.40, 0.45,
             f"TEAM {TEAM.upper()}  ·  PROJECT {PROJECT.upper()}  ·  TRACK  BEST PAYMENT USE CASE  ·  AGENTHON 2026",
             font=MONO, size=10, bold=True, color=CYAN, anchor="middle")

    # Four pain cards
    pains = [
        ("FAILED PAYMENTS",
         "Drop silently. Recovery happens hours late, if at all.",
         RED),
        ("SUSPICIOUS CHARGES",
         "Pile up in inboxes with no consistent triage path.",
         GOLD),
        ("REFUND REQUESTS",
         "Get approved on gut feel, with no audit trail.",
         PURPLE),
        ("DUPLICATE WEBHOOKS",
         "Re-trigger actions and corrupt reconciliation.",
         CYAN),
    ]
    x0 = 0.60
    y = 3.20
    card_w = 2.95
    card_h = 1.70
    gap = 0.10
    for i, (label, body, color) in enumerate(pains):
        x = x0 + i * (card_w + gap)
        add_rect(slide, x, y, card_w, card_h, CARD, line=CARD_EDGE, rounded=True)
        add_rect(slide, x, y, 0.10, card_h, color)
        add_text(slide, x + 0.30, y + 0.20, card_w - 0.50, 0.36, label,
                 font=MONO, size=10, bold=True, color=color)
        add_text(slide, x + 0.30, y + 0.62, card_w - 0.50, card_h - 0.80, body,
                 font=INTER, size=14, color=TEXT)

    # Stat strip
    y2 = 5.15
    stats = [
        ("24/7", "ops loop"),
        ("0", "audit trail today"),
        ("hours", "recovery delay"),
        ("3 humans", "= 3 different decisions"),
    ]
    sw = 2.95
    for i, (big, small) in enumerate(stats):
        x = x0 + i * (sw + gap)
        add_rect(slide, x, y2, sw, 0.95, ROW_BG, line=CARD_EDGE, rounded=True)
        add_text(slide, x, y2 + 0.10, sw, 0.50, big,
                 font=INTER, size=24, bold=True, color=CYAN, align="center")
        add_text(slide, x, y2 + 0.60, sw, 0.30, small,
                 font=MONO, size=9, color=DIM, align="center")

    # Bottom takeaway
    y3 = 6.30
    add_rect(slide, 0.60, y3, 12.10, 0.55, CARD, line=CARD_EDGE, rounded=True)
    add_text(slide, 0.95, y3, 11.40, 0.55,
             "Payment ops in Indonesia is a 24/7 firefight. Merchants need an autonomous operator, not another dashboard.",
             font=INTER, size=13, bold=True, color=GREEN, anchor="middle")


# ---- Slide 2: Solution ---------------------------------------------------

def slide_solution(slide):
    chrome(
        slide,
        eyebrow="02  ·  SOLUTION",
        title="LedgerSoul — an autonomous\npayment-ops agent on DOKU MCP.",
        subtitle="A multi-agent system that observes events, reasons about them, calls DOKU tools, verifies results, and writes an immutable audit trail.",
        page_num=2,
    )

    # Layer separation card (3 cols)
    y = 2.55
    h = 1.55
    cards = [
        ("LAYER 1  ·  PAYMENT RAIL & FRAUD",
         "Owned by banks · networks",
         "3DS · Visa Risk · Mastercard DI · AML",
         RED),
        ("LAYER 2  ·  PAYMENT PLATFORM",
         "Owned by DOKU",
         "Checkout · QRIS · VA · MCP Server",
         PURPLE),
        ("LAYER 3  ·  PAYMENT OPS  (US)",
         "Owned by LedgerSoul",
         "Recovery · triage · refunds · audit",
         CYAN),
    ]
    cw = 3.95
    gap = 0.13
    x0 = 0.60
    for i, (label, owner, body, color) in enumerate(cards):
        x = x0 + i * (cw + gap)
        add_rect(slide, x, y, cw, h, CARD, line=CARD_EDGE, rounded=True)
        add_rect(slide, x, y, 0.10, h, color)
        add_text(slide, x + 0.25, y + 0.18, cw - 0.40, 0.30, label,
                 font=MONO, size=9, bold=True, color=color)
        add_text(slide, x + 0.25, y + 0.50, cw - 0.40, 0.35, owner,
                 font=INTER, size=14, bold=True, color=TEXT)
        add_text(slide, x + 0.25, y + 0.90, cw - 0.40, 0.55, body,
                 font=INTER, size=11, color=DIM)

    # Differentiation row
    y2 = 4.30
    add_text(slide, 0.60, y2, 12.10, 0.30,
             "WHY NOT THE DOKU DASHBOARD OR ZAPIER",
             font=MONO, size=10, bold=True, color=CYAN)
    diffs = [
        ("DOKU DASHBOARD", "Manual clicks. Human judgement. No 24/7 loop.", DIM),
        ("ZAPIER / N8N", "If/then glue. No verification. No audit trace.", DIM),
        ("LEDGERSOUL", "Reasoning + tools + verification + audit, every run.", CYAN),
    ]
    y3 = 4.65
    for i, (label, body, color) in enumerate(diffs):
        x = x0 + i * (cw + gap)
        add_rect(slide, x, y3, cw, 1.05, ROW_BG,
                 line=color if color is CYAN else CARD_EDGE,
                 rounded=True, line_pt=1.0 if color is CYAN else 0.75)
        add_text(slide, x + 0.25, y3 + 0.15, cw - 0.40, 0.28, label,
                 font=MONO, size=9, bold=True, color=color)
        add_text(slide, x + 0.25, y3 + 0.45, cw - 0.40, 0.55, body,
                 font=INTER, size=12,
                 bold=color is CYAN, color=TEXT if color is CYAN else DIM)

    # Bottom: outcome
    y4 = 6.05
    add_rect(slide, 0.60, y4, 12.10, 0.85, CARD, line=CARD_EDGE, rounded=True)
    add_text(slide, 0.95, y4 + 0.10, 11.40, 0.30, "EVERY RUN ENDS IN ONE OF SIX STATES",
             font=MONO, size=9, bold=True, color=CYAN)
    add_text(slide, 0.95, y4 + 0.40, 11.40, 0.40,
             "completed  ·  escalated  ·  duplicate  ·  blocked  ·  failed_verification  ·  error      always with a JSON trace",
             font=INTER, size=13, bold=True, color=GREEN)


# ---- Slide 3: Architecture ----------------------------------------------

def slide_architecture(slide):
    chrome(
        slide,
        eyebrow="03  ·  ARCHITECTURE",
        title="Four agents. One deterministic loop.\nReal DOKU MCP tools.",
        subtitle="Multi-agent reasoning meets a deterministic verification spine — autonomy with proof.",
        page_num=3,
    )

    # Pipeline of 4 agents
    agents = [
        ("REASONING\nAGENT", "LLM",   "classify · extract intent",       PURPLE),
        ("PLANNER\nAGENT",   "rules", "deterministic plan + policy",     CYAN),
        ("EXECUTOR\nAGENT",  "tools", "TOOL_REGISTRY · DOKU MCP",        GOLD),
        ("VERIFIER\nAGENT",  "checks","tool result + audit log",         GREEN),
    ]
    y = 2.55
    h = 1.55
    aw = 2.85
    gap = 0.18
    x0 = 0.60
    for i, (label, badge, body, color) in enumerate(agents):
        x = x0 + i * (aw + gap)
        add_rect(slide, x, y, aw, h, CARD, line=color, rounded=True, line_pt=1.0)
        # badge
        add_rect(slide, x + 0.25, y + 0.18, 0.95, 0.28, BG, line=color, rounded=True)
        add_text(slide, x + 0.25, y + 0.18, 0.95, 0.28, badge.upper(),
                 font=MONO, size=8, bold=True, color=color, align="center", anchor="middle")
        add_text(slide, x + 0.25, y + 0.55, aw - 0.40, 0.65, label,
                 font=INTER, size=15, bold=True, color=TEXT)
        add_text(slide, x + 0.25, y + 1.18, aw - 0.40, 0.30, body,
                 font=MONO, size=9, color=DIM)
        # arrow between cards
        if i < len(agents) - 1:
            ax = x + aw + 0.02
            add_text(slide, ax, y + 0.55, gap - 0.04, 0.45, "▶",
                     font=INTER, size=18, bold=True, color=DIM,
                     align="center", anchor="middle")

    # Tools row
    y2 = 4.35
    add_text(slide, 0.60, y2, 12.10, 0.28,
             "REGISTERED TOOLS  ·  EVERY CALL GOES THROUGH TOOL_REGISTRY",
             font=MONO, size=10, bold=True, color=CYAN)
    tools = [
        "get_transaction_by_invoice_number",
        "create_recovery_link",
        "draft_customer_message",
        "create_approval_request",
        "list_doku_mcp_tools",
        "call_doku_mcp_tool",
        "write_audit_log",
    ]
    chip_y = 4.70
    chip_h = 0.40
    cx = 0.60
    chip_gap = 0.10
    for t in tools:
        cw_chip = 0.18 + 0.085 * len(t)  # rough sizing
        if cx + cw_chip > 12.70:
            cx = 0.60
            chip_y += chip_h + 0.08
        add_rect(slide, cx, chip_y, cw_chip, chip_h, ROW_BG,
                 line=CARD_EDGE, rounded=True)
        add_text(slide, cx, chip_y, cw_chip, chip_h, t,
                 font=MONO, size=10, color=CYAN, align="center", anchor="middle")
        cx += cw_chip + chip_gap

    # Bottom: lifecycle + audit
    y3 = 6.05
    add_rect(slide, 0.60, y3, 12.10, 0.85, CARD, line=CARD_EDGE, rounded=True)
    add_text(slide, 0.95, y3 + 0.10, 11.40, 0.30,
             "12-STAGE AUTONOMOUS LOOP",
             font=MONO, size=9, bold=True, color=CYAN)
    add_text(slide, 0.95, y3 + 0.40, 11.40, 0.40,
             "boot · observe · validate · idempotency · interpret · plan · policy · act · verify · remember · reflect · stop/escalate",
             font=INTER, size=12, color=TEXT)


# ---- Slide 4: Features & Stack ------------------------------------------

def slide_features(slide):
    chrome(
        slide,
        eyebrow="04  ·  FEATURES & STACK",
        title="Built for autonomy, audit,\nand judge-safe demoability.",
        subtitle="Production-shape engineering inside a 12-hour build: tests, lint, judge mode, and real MCP.",
        page_num=4,
    )

    # Two columns — features (left) and stack (right) — sized to leave room
    # for the Judge Mode strip and footer below.
    y = 2.55
    card_h = 3.55
    add_rect(slide, 0.60, y, 6.05, card_h, CARD, line=CARD_EDGE, rounded=True)
    add_text(slide, 0.85, y + 0.18, 5.60, 0.30, "KEY FEATURES",
             font=MONO, size=10, bold=True, color=CYAN)
    features = [
        ("Multi-agent runtime",       "Reasoning + Planner + Executor + Verifier"),
        ("Real DOKU MCP integration", "Sandbox MCP server · 35 tools available"),
        ("Explicit tool registry",    "All tool calls auditable, never arbitrary"),
        ("Deterministic verification","Every result inspected · no fabrication"),
        ("Idempotency & policy",      "Duplicate webhooks · amount thresholds"),
        ("Judge Mode",                "Locked browser demo · redacted traces"),
        ("Audit trail",               "JSON trace + append-only audit_log.jsonl"),
        ("9 evals · 61 tests",        "All passing · ruff clean"),
    ]
    fy = y + 0.52
    row_step = 0.36
    for label, body in features:
        add_rect(slide, 0.85, fy + 0.06, 0.08, 0.18, CYAN)
        add_text(slide, 1.02, fy, 5.40, 0.20, label,
                 font=INTER, size=12, bold=True, color=TEXT)
        add_text(slide, 1.02, fy + 0.18, 5.40, 0.18, body,
                 font=MONO, size=8, color=DIM)
        fy += row_step

    # Right column — stack
    y_r = 2.55
    add_rect(slide, 6.85, y_r, 5.85, card_h, CARD, line=CARD_EDGE, rounded=True)
    add_text(slide, 7.10, y_r + 0.18, 5.40, 0.30, "TECH STACK",
             font=MONO, size=10, bold=True, color=CYAN)
    stack = [
        ("Language",      "Python 3.12"),
        ("Server",        "FastAPI + Uvicorn"),
        ("LLM reasoner",  "Pluggable (OpenAI / Claude) via 9router"),
        ("Payment rails", "DOKU MCP Server (sandbox)"),
        ("Transport",     "HTTP JSON-RPC · MCP 2025-06-18"),
        ("Persistence",   "JSON / JSONL traces · audit log"),
        ("Tests / lint",  "pytest (61 passing) · ruff"),
        ("Deploy",        "Docker · Uvicorn · Judge Mode UI"),
    ]
    sy = y_r + 0.52
    for label, body in stack:
        add_text(slide, 7.10, sy, 1.85, 0.30, label,
                 font=MONO, size=9, color=DIM, anchor="middle")
        add_text(slide, 8.95, sy, 3.65, 0.30, body,
                 font=INTER, size=12, bold=True, color=TEXT, anchor="middle")
        sy += row_step

    # Bottom: judge mode chip — sits BELOW both cards, well above the footer
    y3 = 6.30
    add_rect(slide, 0.60, y3, 12.10, 0.55, CARD, line=CYAN, rounded=True, line_pt=1.0)
    add_text(slide, 0.95, y3, 11.40, 0.55,
             "JUDGE MODE  ·  /judge browser demo  ·  token-protected workflows  ·  redacted traces  ·  /docs · /agent/run · /traces blocked",
             font=MONO, size=10, bold=True, color=CYAN, anchor="middle")


# ---- Slide 5: Future / Impact -------------------------------------------

def slide_future(slide):
    chrome(
        slide,
        eyebrow="05  ·  FUTURE & IMPACT",
        title="From hackathon MVP to the\noperator every DOKU merchant runs.",
        subtitle="Direct impact for Indonesian merchants. Clear roadmap from sandbox to merchant-grade.",
        page_num=5,
    )

    # Roadmap row
    y = 2.55
    h = 1.85
    cw = 3.95
    gap = 0.13
    x0 = 0.60
    phases = [
        ("NOW  ·  v0.1  HACKATHON MVP", CYAN, [
            "12-stage autonomous loop",
            "DOKU MCP sandbox integration",
            "Judge Mode w/ redacted traces",
            "9 evals · 61 tests passing",
        ]),
        ("NEXT  ·  v0.2  PILOT WITH 1 MERCHANT", PURPLE, [
            "Live mode behind explicit flags",
            "Webhook signature validation",
            "Slack / Telegram approvals",
            "SQLite state · richer dashboards",
        ]),
        ("LATER  ·  v1.0  MERCHANT-GRADE", GOLD, [
            "Multi-merchant tenancy",
            "Refund + dispute workflows",
            "Policy editor for ops teams",
            "LLM-assisted planner upgrades",
        ]),
    ]
    for i, (header, color, items) in enumerate(phases):
        x = x0 + i * (cw + gap)
        add_rect(slide, x, y, cw, h, CARD, line=CARD_EDGE, rounded=True)
        add_rect(slide, x, y, 0.10, h, color)
        add_text(slide, x + 0.25, y + 0.18, cw - 0.40, 0.30, header,
                 font=MONO, size=9, bold=True, color=color)
        iy = y + 0.55
        for it in items:
            add_rect(slide, x + 0.30, iy + 0.08, 0.06, 0.16, color)
            add_text(slide, x + 0.45, iy, cw - 0.55, 0.30, it,
                     font=INTER, size=11, color=TEXT)
            iy += 0.30

    # Impact stats
    y2 = 4.65
    impact = [
        ("hours → minutes",  "recovery cycle"),
        ("100%",              "audited runs"),
        ("0",                 "fabricated decisions"),
        ("any merchant",      "with DOKU MCP"),
    ]
    sw = 2.95
    for i, (big, small) in enumerate(impact):
        x = x0 + i * (sw + 0.10)
        add_rect(slide, x, y2, sw, 1.05, ROW_BG, line=CARD_EDGE, rounded=True)
        add_text(slide, x, y2 + 0.15, sw, 0.50, big,
                 font=INTER, size=20, bold=True, color=CYAN, align="center")
        add_text(slide, x, y2 + 0.65, sw, 0.30, small,
                 font=MONO, size=9, color=DIM, align="center")

    # Closing CTA
    y3 = 6.05
    add_rect(slide, 0.60, y3, 12.10, 0.85, CARD, line=CYAN, rounded=True, line_pt=1.0)
    add_text(slide, 0.95, y3 + 0.10, 11.40, 0.30, "WHY DOKU SHOULD CARE",
             font=MONO, size=9, bold=True, color=CYAN)
    add_text(slide, 0.95, y3 + 0.40, 11.40, 0.40,
             "Every DOKU merchant inherits a 24/7 audited operator on day one — without changing checkout, settlement, or fraud rails.",
             font=INTER, size=13, bold=True, color=TEXT)


# ---- Build ---------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    builders = [slide_problem, slide_solution, slide_architecture,
                slide_features, slide_future]
    for build in builders:
        s = prs.slides.add_slide(blank)
        build(s)

    out = Path(f"/home/ubuntu/ledgersoul/OpenClaw2026_{TEAM}_{PROJECT}.pptx")
    prs.save(out)
    print(f"saved {out} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
