"""Insert two new pitch slides into ledgersoul_pitch.pptx.

Adds, after slide 7 (DOKU Integration):
  - Slide 8: "Where We Fit" — layer separation vs bank fraud + DOKU
  - Slide 9: "Why Not the DOKU Dashboard or Zapier" — differentiation
Then renumbers all "NN / 12" footers to "NN / 14".
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# ---- Theme tokens (from existing deck) -----------------------------------
BG = RGBColor(0x07, 0x0B, 0x14)
CYAN = RGBColor(0x4D, 0xE6, 0xF0)
PURPLE = RGBColor(0xA9, 0x7C, 0xFF)
CARD = RGBColor(0x18, 0x24, 0x40)
CARD_EDGE = RGBColor(0x22, 0x31, 0x55)
TEXT = RGBColor(0xE8, 0xEC, 0xF6)
DIM = RGBColor(0x9A, 0xA6, 0xC2)
GREEN = RGBColor(0x69, 0xE6, 0xA3)
RED = RGBColor(0xFF, 0x8A, 0x8A)

INTER = "Inter"
MONO = "JetBrains Mono"


# ---- Helpers --------------------------------------------------------------

def add_rect(slide, x, y, w, h, fill, line=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    if rounded:
        s.adjustments[0] = 0.10
    return s


def add_text(slide, x, y, w, h, text, *, font=INTER, size=14, bold=False,
             color=TEXT, align="left", anchor="top"):
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def slide_chrome(slide, eyebrow, title, subtitle, page_num, total):
    add_rect(slide, 0, 0, 13.33, 7.5, BG)
    # Eyebrow accent square + text
    add_rect(slide, 0.60, 0.55, 0.18, 0.18, CYAN)
    add_text(slide, 0.85, 0.50, 8.0, 0.30, eyebrow,
             font=MONO, size=10, bold=True, color=CYAN)
    # Title — give two-line headlines enough vertical room so subtitle
    # never collides. Two lines of 30pt @ 1.15 line-height ≈ 0.95 in.
    add_text(slide, 0.60, 0.85, 12.10, 1.10, title,
             font=INTER, size=30, bold=True, color=TEXT)
    # Subtitle sits well below the title block
    if subtitle:
        add_text(slide, 0.60, 2.00, 11.50, 0.40, subtitle,
                 font=INTER, size=14, color=DIM)
    # Footer
    add_text(slide, 0.60, 7.05, 9.0, 0.30,
             "LEDGERSOUL  ·  AUTONOMOUS PAYMENT OPERATIONS AGENT  ·  DOKU HACKATHON",
             font=MONO, size=8, bold=True, color=DIM)
    add_text(slide, 11.60, 7.05, 1.20, 0.30,
             f"{page_num:02d} / {total:02d}",
             font=MONO, size=8, bold=True, color=DIM, align="right")


# ---- Slide builders -------------------------------------------------------

def build_where_we_fit(slide):
    slide_chrome(
        slide,
        eyebrow="07a  ·  WHERE WE FIT",
        title="We are not a payment gateway,\nand not a fraud system.",
        subtitle="LedgerSoul lives in the operational layer above DOKU — after the bank decided, after the gateway settled.",
        page_num=8, total=14,
    )

    # 3 stacked horizontal layer bars
    layers = [
        # (color, eyebrow, label, examples, owner_tag, owner_color)
        (RED, "LAYER 1  ·  PAYMENT RAIL & FRAUD",
         "Bank, network, and acquirer-level controls.",
         "3DS · velocity rules · Visa Risk · Mastercard DI · AML / sanctions screening",
         "OWNED BY  BANKS · NETWORKS",
         RED),
        (PURPLE, "LAYER 2  ·  PAYMENT PLATFORM",
         "Gateway, checkout, settlement, dashboard.",
         "DOKU Checkout · QRIS · Virtual Account · Payment Link · Merchant Dashboard · DOKU MCP",
         "OWNED BY  DOKU",
         PURPLE),
        (CYAN, "LAYER 3  ·  PAYMENT OPS  (US)",
         "What happens after the event lands.",
         "Reconciliation · failed-payment recovery · refund triage · escalation · audit trail",
         "OWNED BY  LEDGERSOUL",
         CYAN),
    ]

    top = 2.55
    bar_h = 1.15
    gap = 0.15
    for i, (accent, eye, label, examples, owner, owner_color) in enumerate(layers):
        y = top + i * (bar_h + gap)
        # Card
        add_rect(slide, 0.60, y, 12.10, bar_h, CARD, line=CARD_EDGE, rounded=True)
        # Left accent rail
        add_rect(slide, 0.60, y, 0.10, bar_h, accent)
        # Eyebrow
        add_text(slide, 0.95, y + 0.15, 6.5, 0.28, eye,
                 font=MONO, size=10, bold=True, color=accent)
        # Label (bold)
        add_text(slide, 0.95, y + 0.40, 7.5, 0.36, label,
                 font=INTER, size=16, bold=True, color=TEXT)
        # Examples
        add_text(slide, 0.95, y + 0.76, 8.5, 0.34, examples,
                 font=INTER, size=11, color=DIM)
        # Owner tag (right side)
        add_rect(slide, 10.30, y + 0.42, 2.20, 0.36, BG, line=owner_color, rounded=True)
        add_text(slide, 10.30, y + 0.42, 2.20, 0.36, owner,
                 font=MONO, size=8, bold=True, color=owner_color, align="center", anchor="middle")

    # Bottom takeaway strip
    y = 6.40
    add_rect(slide, 0.60, y, 12.10, 0.50, CARD, line=CARD_EDGE, rounded=True)
    add_text(slide, 0.95, y, 11.40, 0.50,
             "We complement fraud detection. We extend the DOKU dashboard. We don't compete with either.",
             font=INTER, size=13, bold=True, color=GREEN, anchor="middle")


def build_differentiation(slide):
    slide_chrome(
        slide,
        eyebrow="07b  ·  WHY LEDGERSOUL",
        title="Why not just the DOKU dashboard\nor a Zapier workflow?",
        subtitle="The dashboard is manual. Zapier is brittle glue. LedgerSoul is an auditable autonomous operator.",
        page_num=9, total=14,
    )

    # 4-column comparison table
    cols = [
        ("CAPABILITY",        None,   None),
        ("DOKU DASHBOARD",    DIM,    "manual"),
        ("ZAPIER / N8N",      DIM,    "glue"),
        ("LEDGERSOUL",        CYAN,   "agent"),
    ]
    rows = [
        ("Triggered by events",            "human clicks",       "single trigger",     "any payment event"),
        ("Decision making",                "human judgement",    "if/then rules",      "planner + policy"),
        ("Verifies tool results",          "—",                  "—",                  "yes, every run"),
        ("Idempotency on duplicates",      "manual",             "fragile",            "deterministic"),
        ("Human-approval thresholds",      "manual",             "custom code",        "built-in policy"),
        ("Audit trace per run",            "partial logs",       "run history",        "JSON trace + audit"),
        ("Identity & contract",            "—",                  "—",                  "soul.md + agent.md"),
        ("Operates 24/7 unattended",       "no",                 "yes (brittle)",      "yes (deterministic)"),
    ]

    x0 = 0.60
    col_widths = [3.20, 2.80, 2.80, 3.30]   # sum = 12.10
    y0 = 2.55
    row_h = 0.36
    header_h = 0.44
    table_total_h = header_h + len(rows) * row_h

    # LedgerSoul column highlight strip — sized exactly to table height
    ls_x = x0 + col_widths[0] + col_widths[1] + col_widths[2]
    ls_w = col_widths[3]
    add_rect(slide, ls_x, y0, ls_w, table_total_h,
             RGBColor(0x10, 0x1B, 0x33), line=CYAN, rounded=True)

    # Header band on the three left columns only
    add_rect(slide, x0, y0, col_widths[0] + col_widths[1] + col_widths[2],
             header_h, CARD, line=CARD_EDGE, rounded=True)

    # Column headers
    cx = x0
    for (label, color, _), w in zip(cols, col_widths):
        if label == "LEDGERSOUL":
            c = CYAN
        else:
            c = DIM
        add_text(slide, cx + 0.20, y0, w - 0.40, header_h, label,
                 font=MONO, size=10, bold=True, color=c, anchor="middle")
        cx += w

    # Body rows — zebra on left three columns
    for i, row in enumerate(rows):
        ry = y0 + header_h + i * row_h
        if i % 2 == 0:
            add_rect(slide, x0, ry,
                     col_widths[0] + col_widths[1] + col_widths[2],
                     row_h, RGBColor(0x0D, 0x14, 0x28))
        cx = x0
        for j, (text_, w) in enumerate(zip(row, col_widths)):
            if j == 0:
                color = TEXT
                bold = True
            elif j == 3:
                color = CYAN
                bold = True
            else:
                color = DIM
                bold = False
            add_text(slide, cx + 0.20, ry, w - 0.40, row_h, text_,
                     font=INTER, size=11, bold=bold, color=color, anchor="middle")
            cx += w

    # Bottom callout — pain points solved
    y = 6.40
    add_rect(slide, 0.60, y, 12.10, 0.50, CARD, line=CARD_EDGE, rounded=True)
    add_text(slide, 0.95, y + 0.05, 11.40, 0.18, "PAIN POINTS WE SOLVE",
             font=MONO, size=8, bold=True, color=CYAN)
    add_text(slide, 0.95, y + 0.22, 11.40, 0.26,
             "Silent failed payments  ·  inconsistent triage  ·  unauditable refunds  ·  duplicate webhooks  ·  ops staff stuck firefighting in tabs",
             font=INTER, size=11, color=TEXT)


# ---- Slide reorder + footer renumber -------------------------------------

def move_slide(prs, old_index: int, new_index: int) -> None:
    """Reorder a slide inside the presentation's <p:sldIdLst>."""
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    moved = sldIds[old_index]
    sldIdLst.remove(moved)
    # Insert at the new position
    if new_index >= len(sldIds) - 1:
        sldIdLst.append(moved)
    else:
        ref = sldIds[new_index] if new_index < old_index else sldIds[new_index + 1]
        # Re-fetch updated list and find ref again because `moved` was removed
        current = list(sldIdLst)
        # Compute target position
        if new_index >= len(current):
            sldIdLst.append(moved)
        else:
            sldIdLst.insert(sldIdLst.index(current[new_index]), moved)


def renumber_footers(prs, total: int) -> None:
    """Replace "NN / 12" footers with "NN / total" using new slide order."""
    import re
    for idx, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    m = re.match(r"^\s*(\d{1,2})\s*/\s*\d{1,2}\s*$", run.text)
                    if m:
                        run.text = f"{idx:02d} / {total:02d}"


def main() -> None:
    src = Path("/home/ubuntu/ledgersoul/ledgersoul_pitch.pptx")
    prs = Presentation(src)

    blank = prs.slide_layouts[6]  # Blank
    # Append two new slides at the end first
    s_fit = prs.slides.add_slide(blank)
    build_where_we_fit(s_fit)
    s_diff = prs.slides.add_slide(blank)
    build_differentiation(s_diff)

    # Move them: they were appended at indices [12, 13] (0-based: 12, 13).
    # We want them at positions 7 and 8 (0-based) — i.e. directly after slide
    # index 6 (DOKU Integration is slide 7 / index 6).
    # First move the differentiation slide (last) to index 7 — it will end up
    # right after DOKU Integration.
    move_slide(prs, old_index=13, new_index=7)
    # Now the where-we-fit slide is at index 12 (it shifted down by one). Move
    # it to index 7 so it sits before differentiation.
    # After the previous move: order tail is ..., DOKU(6), Diff(7), JudgeMode(8), ..., WhereWeFit(13)
    move_slide(prs, old_index=13, new_index=7)

    # Renumber footers (00 / 14)
    renumber_footers(prs, total=14)

    out = Path("/home/ubuntu/ledgersoul/ledgersoul_pitch.pptx")
    prs.save(out)
    print(f"saved {out} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
